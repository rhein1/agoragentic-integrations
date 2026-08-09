#!/usr/bin/env python3
"""Generate and independently verify public-safe semantic fixtures."""

from __future__ import annotations

import csv
import hashlib
from importlib.metadata import version
from io import BytesIO
import json
from pathlib import Path
import shutil

from docx import Document
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

ROOT = Path(__file__).resolve().parent
GENERATED = ROOT / "generated"
MANIFEST = GENERATED / "fixture-manifest.json"
REQUIREMENTS_LOCK = ROOT / "requirements.txt"
GENERATOR_DISTRIBUTIONS = (
    "charset-normalizer",
    "et-xmlfile",
    "lxml",
    "openpyxl",
    "Pillow",
    "python-docx",
    "python-pptx",
    "reportlab",
    "typing-extensions",
    "xlsxwriter",
)


def sha256_file(path: Path) -> str:
    return f"sha256:{hashlib.sha256(path.read_bytes()).hexdigest()}"


def write_csv() -> None:
    path = GENERATED / "baseline.csv"
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["name", "value"])
        writer.writerow(["alpha", 1])
        writer.writerow(["beta", 2])


def write_long_csv() -> None:
    path = GENERATED / "long-content.csv"
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["payload"])
        writer.writerow([f"{'A' * 2_400}FINAL_MARKER"])


def write_xlsx() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data Sheet"

    sheet["A1"] = "Visible row"
    sheet["A2"] = "Hidden row"
    sheet.row_dimensions[2].hidden = True

    sheet["B1"] = "Hidden column"
    sheet.column_dimensions["B"].hidden = True

    sheet["C1"] = "Displayed percent"
    sheet["C2"] = 0.075
    sheet["C2"].number_format = "0.0%"

    sheet["D1"] = "Formula"
    sheet["D2"] = "=C2"

    sheet.merge_cells("F1:O3")
    sheet["F1"] = "Merged heading"

    xlsx_path = GENERATED / "semantic-risks.xlsx"
    workbook.save(xlsx_path)
    shutil.copyfile(xlsx_path, GENERATED / "semantic-risks.xlsm")


def write_docx() -> None:
    document = Document()
    document.add_heading("Nested table evidence", level=1)
    outer = document.add_table(rows=1, cols=1)
    cell = outer.cell(0, 0)
    cell.paragraphs[0].add_run("Section A")
    inner = cell.add_table(rows=3, cols=2)
    inner.cell(0, 0).text = "Metric"
    inner.cell(0, 1).text = "Value"
    inner.cell(1, 0).text = "Height"
    inner.cell(1, 1).text = "120"
    inner.cell(2, 0).text = "Weight"
    inner.cell(2, 1).text = "34"
    document.save(GENERATED / "nested-table.docx")


def add_textbox(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = text


def write_pptx() -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Titled slide"

    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(second, "Second slide, no title placeholder.")

    third = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(third, "Third slide, also untitled.")

    presentation.save(GENERATED / "untitled-slides.pptx")


def write_image_only_pdf() -> None:
    image = Image.new("RGB", (300, 120), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((20, 20, 280, 100), outline="black", width=3)
    draw.text((40, 50), "pixels only", fill="black")

    png = BytesIO()
    image.save(png, format="PNG")
    png.seek(0)

    path = GENERATED / "image-only.pdf"
    pdf = canvas.Canvas(str(path), pagesize=(360, 180))
    pdf.drawImage(ImageReader(png), 30, 30, width=300, height=120)
    pdf.showPage()
    pdf.save()


def verify_sources() -> dict[str, list[str]]:
    with (GENERATED / "baseline.csv").open(encoding="utf-8", newline="") as handle:
        assert list(csv.reader(handle)) == [["name", "value"], ["alpha", "1"], ["beta", "2"]]

    with (GENERATED / "long-content.csv").open(encoding="utf-8", newline="") as handle:
        rows = list(csv.reader(handle))
    assert rows[0] == ["payload"]
    assert len(rows[1][0]) == 2_412 and rows[1][0].endswith("FINAL_MARKER")

    workbook = load_workbook(GENERATED / "semantic-risks.xlsx", data_only=False)
    sheet = workbook["Data Sheet"]
    assert sheet.row_dimensions[2].hidden is True
    assert sheet.column_dimensions["B"].hidden is True
    assert sheet["C2"].value == 0.075 and sheet["C2"].number_format == "0.0%"
    assert sheet["D2"].value == "=C2"
    assert "F1:O3" in {str(item) for item in sheet.merged_cells.ranges}
    assert (GENERATED / "semantic-risks.xlsx").read_bytes() == (GENERATED / "semantic-risks.xlsm").read_bytes()

    document = Document(GENERATED / "nested-table.docx")
    nested = document.tables[0].cell(0, 0).tables[0]
    assert [[cell.text for cell in row.cells] for row in nested.rows] == [
        ["Metric", "Value"],
        ["Height", "120"],
        ["Weight", "34"],
    ]

    presentation = Presentation(GENERATED / "untitled-slides.pptx")
    assert len(presentation.slides) == 3
    assert presentation.slides[0].shapes.title.text == "Titled slide"
    assert presentation.slides[1].shapes.title is None
    assert presentation.slides[2].shapes.title is None
    slide_text = [
        " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))
        for slide in presentation.slides
    ]
    assert "Second slide, no title placeholder." in slide_text[1]
    assert "Third slide, also untitled." in slide_text[2]

    pdf_bytes = (GENERATED / "image-only.pdf").read_bytes()
    assert pdf_bytes.startswith(b"%PDF-") and b"/Subtype /Image" in pdf_bytes

    return {
        "baseline.csv": ["exact_csv_rows_verified"],
        "long-content.csv": ["long_cell_length_and_terminal_marker_verified"],
        "semantic-risks.xlsx": [
            "hidden_row_verified",
            "hidden_column_verified",
            "display_format_verified",
            "formula_verified",
            "merged_range_verified",
        ],
        "semantic-risks.xlsm": ["byte_identical_alias_fixture_verified"],
        "nested-table.docx": ["nested_table_cells_verified"],
        "untitled-slides.pptx": ["untitled_slide_structure_and_text_verified"],
        "image-only.pdf": ["pdf_contains_embedded_image_object"],
    }


def write_manifest(facts: dict[str, list[str]]) -> None:
    fixtures = {}
    for filename, verified_facts in sorted(facts.items()):
        path = GENERATED / filename
        fixtures[filename] = {
            "size_bytes": path.stat().st_size,
            "source_hash": sha256_file(path),
            "verified_source_facts": verified_facts,
        }
    manifest = {
        "schema": "agoragentic.anydoc-conformance-fixtures.v1",
        "generator_dependencies": {
            distribution: version(distribution)
            for distribution in GENERATOR_DISTRIBUTIONS
        },
        "requirements_lock": {
            "path": "requirements.txt",
            "source_hash": sha256_file(REQUIREMENTS_LOCK),
            "hashes_required": True,
        },
        "fixtures": fixtures,
    }
    MANIFEST.write_text(f"{json.dumps(manifest, indent=2)}\n", encoding="utf-8")


def main() -> None:
    if GENERATED.exists():
        shutil.rmtree(GENERATED)
    GENERATED.mkdir(parents=True)
    write_csv()
    write_long_csv()
    write_xlsx()
    write_docx()
    write_pptx()
    write_image_only_pdf()
    facts = verify_sources()
    write_manifest(facts)
    for path in sorted(GENERATED.iterdir()):
        print(f"generated {path.name} ({path.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
